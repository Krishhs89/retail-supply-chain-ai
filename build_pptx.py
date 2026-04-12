"""
Generates presentation.pptx for the Retail Supply Chain Optimization AI project.
Run: python build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.oxml.ns as nsmap
from lxml import etree
import copy

# ── Colours ──────────────────────────────────────────────────────────────────
BLUE      = RGBColor(0x00, 0x71, 0xCE)   # Walmart blue
DARKBLUE  = RGBColor(0x00, 0x4A, 0x8F)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
DARK      = RGBColor(0x1A, 0x1A, 0x2E)
GREY      = RGBColor(0x6C, 0x75, 0x7D)
LIGHTGREY = RGBColor(0xF8, 0xF9, 0xFA)
GREEN     = RGBColor(0x28, 0xA7, 0x45)
AMBER     = RGBColor(0xF3, 0x9C, 0x12)
RED       = RGBColor(0xDC, 0x35, 0x45)
ROWALT    = RGBColor(0xF0, 0xF4, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # completely blank


# ─── Helpers ─────────────────────────────────────────────────────────────────

def add_rect(slide, left, top, width, height, fill=None, line=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=14,
                bold=False, color=DARK, align=PP_ALIGN.LEFT,
                italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_para(tf, text, font_size=12, bold=False, color=DARK,
             align=PP_ALIGN.LEFT, space_before=0, bullet=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    if bullet:
        p.level = 1
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def slide_header(slide, title, subtitle=None, bg=BLUE):
    """Top banner with title."""
    banner = add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), fill=bg)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(12.3), Inches(0.7),
                title, font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.82), Inches(12.3), Inches(0.38),
                    subtitle, font_size=13, color=RGBColor(0xCC, 0xE5, 0xFF))


def table_slide(slide, headers, rows, top=Inches(1.4), left=Inches(0.4),
                width=Inches(12.5), row_height=Inches(0.38)):
    col_count = len(headers)
    col_w = width / col_count

    # Header row
    for ci, h in enumerate(headers):
        rect = add_rect(slide, left + col_w * ci, top, col_w, row_height, fill=BLUE)
        add_textbox(slide, left + col_w * ci + Inches(0.05), top + Inches(0.05),
                    col_w - Inches(0.1), row_height - Inches(0.05),
                    h, font_size=11, bold=True, color=WHITE)

    # Data rows
    for ri, row in enumerate(rows):
        bg = ROWALT if ri % 2 == 0 else WHITE
        for ci, cell in enumerate(row):
            rect = add_rect(slide, left + col_w * ci, top + row_height * (ri + 1),
                            col_w, row_height, fill=bg,
                            line=RGBColor(0xDE, 0xE2, 0xE6))
            add_textbox(slide, left + col_w * ci + Inches(0.05),
                        top + row_height * (ri + 1) + Inches(0.04),
                        col_w - Inches(0.1), row_height - Inches(0.06),
                        cell, font_size=10, color=DARK)


def bullet_box(slide, items, left, top, width, height,
               icon="•", font_size=13, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{icon}  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)

# Full blue background
add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=BLUE)
# White bottom strip
add_rect(sl, 0, Inches(5.8), SLIDE_W, Inches(1.7), fill=DARKBLUE)

add_textbox(sl, Inches(1.0), Inches(1.2), Inches(11.3), Inches(1.2),
            "🏪  Retail Supply Chain Optimization AI",
            font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.0), Inches(2.5), Inches(11.3), Inches(0.6),
            "Multi-Agent Agentic AI System for Enterprise Retail Decision-Making",
            font_size=18, color=RGBColor(0xCC, 0xE5, 0xFF), align=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.0), Inches(3.3), Inches(11.3), Inches(0.5),
            "Powered by Claude claude-sonnet-4-6  ·  LangGraph  ·  Streamlit",
            font_size=13, color=RGBColor(0xAD, 0xD8, 0xE6), align=PP_ALIGN.CENTER, italic=True)

add_textbox(sl, Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.5),
            "Interview Kickstart — Agentic AI Capstone Project  ·  2026",
            font_size=12, color=RGBColor(0xAD, 0xD8, 0xE6), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "The Problem: Decisions Don't Live in Isolation",
             "Every supply chain decision cascades across 5 domains simultaneously")

# Left column
add_rect(sl, Inches(0.4), Inches(1.35), Inches(6.0), Inches(4.8), fill=RGBColor(0xF0, 0xF4, 0xFF),
         line=RGBColor(0x00, 0x71, 0xCE))
add_textbox(sl, Inches(0.55), Inches(1.45), Inches(5.7), Inches(0.4),
            "A 10% diaper price increase triggers:", font_size=13, bold=True, color=BLUE)
items_l = [
    "📉  -14% demand volume (elasticity = -1.4)",
    "📦  Replenishment PO recalculation across 4 DCs",
    "🚛  Carrier load change in SE and MW regions",
    "💰  Margin shift after vendor trade dollar netting",
    "⚠️   Potential conflict with active promotions",
]
bullet_box(sl, items_l, Inches(0.55), Inches(1.9), Inches(5.7), Inches(3.4),
           icon="", font_size=12)

# Right column
add_rect(sl, Inches(6.9), Inches(1.35), Inches(6.0), Inches(4.8), fill=RGBColor(0xFF, 0xF3, 0xCD),
         line=AMBER)
add_textbox(sl, Inches(7.05), Inches(1.45), Inches(5.7), Inches(0.4),
            "Why this matters:", font_size=13, bold=True, color=RGBColor(0x85, 0x64, 0x04))
items_r = [
    "Human expert: 2–4 hours per decision",
    "This system: seconds, fully connected",
    "Pricing ↔ Demand ↔ Inventory",
    "Supply Chain ↔ Finance — all linked",
    "$1B+ decisions made daily in retail",
]
bullet_box(sl, items_r, Inches(7.05), Inches(1.9), Inches(5.7), Inches(3.0),
           icon="→", font_size=12, color=RGBColor(0x4A, 0x3C, 0x00))

# Bottom stat bar
add_rect(sl, Inches(0.4), Inches(6.3), Inches(12.5), Inches(0.7), fill=BLUE)
add_textbox(sl, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.5),
            "$1B+ revenue decisions in retail require connected, real-time AI reasoning — not spreadsheets.",
            font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — WHAT WE BUILT
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "What We Built", "Two pipelines · 17 tools · 11-tab Streamlit UI · Deployed on Streamlit Cloud")

# V1 box
add_rect(sl, Inches(0.4), Inches(1.4), Inches(5.8), Inches(2.6),
         fill=RGBColor(0xCC, 0xE5, 0xFF), line=BLUE)
add_textbox(sl, Inches(0.55), Inches(1.5), Inches(5.5), Inches(0.4),
            "V1 — Agentic Loop", font_size=14, bold=True, color=BLUE)
bullet_box(sl, [
    "Single Claude claude-sonnet-4-6 agent",
    "All 17 tools available simultaneously",
    "Adaptive iterations: 6 / 10 / 20",
    "History summarization at 12 messages",
    "Configurable via sidebar slider",
], Inches(0.55), Inches(1.95), Inches(5.5), Inches(2.0), icon="✓", font_size=11)

# V2 box
add_rect(sl, Inches(6.55), Inches(1.4), Inches(6.35), Inches(2.6),
         fill=RGBColor(0xD4, 0xED, 0xDA), line=GREEN)
add_textbox(sl, Inches(6.7), Inches(1.5), Inches(6.05), Inches(0.4),
            "V2 — LangGraph Multi-Agent Graph", font_size=14, bold=True, color=GREEN)
bullet_box(sl, [
    "12 specialist nodes in directed graph",
    "Router → Domain Nodes → Synthesizer",
    "Each node: curated tools + focused prompt",
    "Node-by-node execution trace",
    "Conditional edges per intent type",
], Inches(6.7), Inches(1.95), Inches(6.05), Inches(2.0), icon="✓", font_size=11)

# Tools row
add_textbox(sl, Inches(0.4), Inches(4.15), Inches(12.5), Inches(0.35),
            "17 Operational Tools", font_size=14, bold=True, color=BLUE)

tool_cats = [
    ("Pricing", "simulate_price_change\nget_competitive_pricing\nadjust_promotional_price"),
    ("Demand", "get_demand_forecast\nget_forecast_accuracy\nanalyze_demand_variables"),
    ("Inventory", "get_inventory_levels\ncalculate_stockout_risk\nget_reorder_recommendations"),
    ("Supply Chain", "get_carrier_status\nfind_alternate_carriers\nget_dc_inventory"),
    ("Finance & Scenario", "calculate_revenue_impact\ncalculate_carrying_cost\ndetect_scenario_conflicts\ncompare_scenarios\nget_replenishment_schedule"),
]
cat_colors = [RGBColor(0xCC, 0xE5, 0xFF), RGBColor(0xD4, 0xED, 0xDA),
              RGBColor(0xFF, 0xF3, 0xCD), RGBColor(0xF8, 0xD7, 0xDA), RGBColor(0xE8, 0xD5, 0xFF)]
cat_w = Inches(2.42)
for i, (cat, tools) in enumerate(tool_cats):
    cx = Inches(0.4) + cat_w * i
    add_rect(sl, cx, Inches(4.55), cat_w - Inches(0.08), Inches(2.55),
             fill=cat_colors[i], line=RGBColor(0xDE, 0xE2, 0xE6))
    add_textbox(sl, cx + Inches(0.08), Inches(4.6), cat_w - Inches(0.18), Inches(0.32),
                cat, font_size=10, bold=True, color=DARK)
    add_textbox(sl, cx + Inches(0.08), Inches(4.95), cat_w - Inches(0.18), Inches(2.0),
                tools, font_size=9, color=GREY)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "System Architecture", "LangGraph V2 routing topology")

arch_text = """                    USER QUERY
                         │
                         ▼
               ┌──────────────────┐
               │     ROUTER       │  classifies intent, extracts SKU + entities
               └──────┬───────────┘
                      │
      ┌───────────┬───┴────────┬───────────┬──────────────┐
      ▼           ▼            ▼           ▼              ▼
  Price        Supply      Demand      Scenario      Shelf
  Cascade    Disruption   Forecast    Planning    Replenishment
      │           │            │                       │
      ▼           ▼            ▼                       ▼
  Inventory   Carrier      Accuracy               Perishable
  Node        Node         Gate                   Check
      │           │            │                       │
      └─────────┐ │ ┌──────────┘                       │
                ▼ ▼ ▼                                   │
           Financial                                    │
           Impact                                       │
                └───────────────┬───────────────────────┘
                                ▼
                       ┌────────────────┐
                       │  SYNTHESIZER   │  merges all node outputs
                       └────────┬───────┘
                                ▼
                               END"""

txBox = sl.shapes.add_textbox(Inches(0.5), Inches(1.35), Inches(12.3), Inches(5.8))
tf = txBox.text_frame
tf.word_wrap = False
p = tf.paragraphs[0]
run = p.add_run()
run.text = arch_text
run.font.size = Pt(9.5)
run.font.name = "Courier New"
run.font.color.rgb = DARK


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — EDGE CASES TABLE
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "17 Production-Grade Edge Cases Handled")

table_slide(sl,
    headers=["Edge Case", "Behavior", "Impact"],
    rows=[
        ["Asymmetric price elasticity", "Tobacco/diapers recover 40–70% on price cuts", "Prevents over-optimistic promo ROI"],
        ["Perishable cap", "Dairy orders hard-capped at 3 days supply", "Eliminates markdown waste"],
        ["Regional carrier gap", "SE has only TruckCo D for diapers (+45%)", "Surfaces hidden cost risk"],
        ["Forecast accuracy gate", "<60% accuracy blocks PO system", "Prevents bad-data-driven orders"],
        ["CI widening", "+15% per 4-week horizon period", "Honest uncertainty at 8 weeks: ±30%"],
        ["VMI exclusion", "Vendor-managed inventory excl. from carrying cost", "Accurate cost reporting"],
        ["Vendor trade netting", "Subsidies netted before margin reporting", "True promo profitability"],
        ["Replenishment lag", "3–4 days + 30% chance of +3 extra days", "Realistic lead time planning"],
        ["Scenario conflict", "Promo + supply disruption = CRITICAL flag", "Prevents simultaneous blunders"],
        ["Data provenance", "OLAP (24h) vs WMS (15min) freshness warning", "Right data for right decision"],
    ],
    top=Inches(1.35),
    row_height=Inches(0.505),
)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — SCENARIO: CARRIER STRIKE + PROMO
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "Live Scenario: Carrier Strike + Promotion Conflict",
             "The system detects CRITICAL conflict and recommends deferring the promotion")

# Scenario box
add_rect(sl, Inches(0.4), Inches(1.35), Inches(12.5), Inches(0.8),
         fill=RGBColor(0xFF, 0xF3, 0xCD), line=AMBER)
add_textbox(sl, Inches(0.55), Inches(1.42), Inches(12.2), Inches(0.65),
            "⚠️  TruckCo B (diapers, SE + MW) on strike · Category manager wants 10% promo price cut on HUG48-3",
            font_size=13, bold=True, color=RGBColor(0x85, 0x64, 0x04))

steps = [
    ("1. Conflict Detection", "CRITICAL: Promo creates +14% demand surge while supply is constrained", RED),
    ("2. Regional Gap", "SE has NO diaper carrier except TruckCo D — at +45% cost premium", RED),
    ("3. Days-of-Supply", "DC-SE hits safety stock threshold in 4 days without alternate carrier", AMBER),
    ("4. Financial Impact", "Revenue-at-risk: ~$890K over 14-day strike window", AMBER),
    ("5. Recommendation", "Defer promotion until carrier restored OR pre-build DC-SE inventory now", GREEN),
]

for i, (step, detail, color) in enumerate(steps):
    top = Inches(2.3) + Inches(0.88) * i
    add_rect(sl, Inches(0.4), top, Inches(2.8), Inches(0.72), fill=color,
             line=RGBColor(0xDE, 0xE2, 0xE6))
    add_textbox(sl, Inches(0.5), top + Inches(0.12), Inches(2.6), Inches(0.5),
                step, font_size=11, bold=True, color=WHITE)
    add_rect(sl, Inches(3.3), top, Inches(9.65), Inches(0.72),
             fill=LIGHTGREY, line=RGBColor(0xDE, 0xE2, 0xE6))
    add_textbox(sl, Inches(3.45), top + Inches(0.12), Inches(9.4), Inches(0.5),
                detail, font_size=12, color=DARK)

add_rect(sl, Inches(0.4), Inches(6.75), Inches(12.5), Inches(0.55), fill=RED)
add_textbox(sl, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
            "❌  Proceed with promotion = stockout risk in SE within 4 days + 45% carrier premium eating all promo margin",
            font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SCENARIO: PRICE CASCADE
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "Live Scenario: Full Price Cascade Analysis",
             'Query: "Raise HUG48-3 diaper price from $12.99 to $14.49 — simulate full cascade"')

table_slide(sl,
    headers=["Step", "What Happens", "Numbers"],
    rows=[
        ["Price change", "Unit price increases +$1.50", "$12.99 → $14.49  (+11.5%)"],
        ["Demand impact", "Elasticity = -1.4 applied to price increase", "-16.1% volume drop"],
        ["Asymmetric check", "Recovery factor = 0.70 logged for future reversal", "Applies only to price cuts"],
        ["PO recalculation", "Lower demand signal → smaller replenishment orders", "-16% PO quantity across DCs"],
        ["Margin impact", "Higher price × lower volume — net effect calculated", "+4.2% net revenue"],
        ["Vendor trade netting", "8% vendor subsidy on promos only — not triggered here", "No promo subsidy applied"],
        ["Carrying cost", "Fewer owned units reduces inventory holding cost", "-$12K annual carrying cost"],
        ["Recommendation", "Price increase accretive — carrier capacity headroom exists", "✅ Recommend proceed"],
    ],
    top=Inches(1.35),
    row_height=Inches(0.63),
)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — UI OVERVIEW
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "11-Tab Streamlit Interface", "Every tab is powered by the same mock data and tool layer")

tab_data = [
    ("Dashboard", "Live carrier status chips, DC inventory KPIs, strike timeline", "📊"),
    ("Chat", "Streaming V1/V2 agent with live tool tracking + node trace", "💬"),
    ("Price Cascade", "Price change → demand → PO → financial waterfall chart", "🔺"),
    ("Supply Alert", "Carrier strike → stockout bars → alternate carrier options", "🚛"),
    ("Demand Forecast", "8-week fan chart with CI widening + 15 variable contributions", "📈"),
    ("Scenario Planner", "3-scenario side-by-side comparison + conflict detection", "⚖️"),
    ("Shelf & Store", "Replenishment Gantt + perishable cap + planogram check", "🏬"),
    ("Financial Impact", "P&L waterfall: revenue → margin → trade → VMI → net", "💰"),
    ("Data Sources", "Provenance tracker — WMS vs OLAP freshness comparison", "🗄️"),
    ("Workflow", "Step-by-step integrated scenario builder + ripple timeline (NEW)", "🔄"),
    ("Flow Map", "LangGraph DAG with live green/gray execution trace (NEW)", "🗺️"),
]

col_w = Inches(4.0)
row_h = Inches(0.545)
for i, (tab, desc, icon) in enumerate(tab_data):
    col = i % 3
    row = i // 3
    left = Inches(0.4) + col_w * col
    top  = Inches(1.38) + row_h * row
    bg = RGBColor(0xCC, 0xE5, 0xFF) if i >= 9 else LIGHTGREY
    add_rect(sl, left, top, col_w - Inches(0.1), row_h - Inches(0.05),
             fill=bg, line=RGBColor(0xDE, 0xE2, 0xE6))
    add_textbox(sl, left + Inches(0.1), top + Inches(0.06),
                col_w - Inches(0.2), row_h - Inches(0.1),
                f"{icon}  {tab}  —  {desc}", font_size=10, color=DARK)

add_textbox(sl, Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.25),
            "Blue tabs = new in this release", font_size=10, italic=True, color=BLUE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — RESULTS
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "What Was Achieved", "Technical + business capabilities delivered")

tech = [
    "Full multi-agent agentic loop (V1)",
    "LangGraph 12-node directed graph (V2)",
    "17 domain tools with realistic mock data",
    "Asymmetric elasticity modeling",
    "Scenario conflict detection engine",
    "Data provenance tracking (OLAP/WMS/OLTP)",
    "Context window management (12-msg threshold)",
    "Streaming UI with live tool tracking",
    "Deployed on Streamlit Community Cloud",
    "GitHub repository with full documentation",
]

biz = [
    "Price → demand → PO → financial cascade",
    "Regional carrier gap detection",
    "Perishable inventory capping (3-day hard cap)",
    "Vendor trade dollar netting",
    "VMI exclusion from carrying costs",
    "Forecast accuracy gating (60% floor)",
    "Multi-scenario comparison with ranking",
    "Replenishment chain lag modeling",
    "Planogram overflow handling",
    "GUIDE.md specialist decision-maker reference",
]

# Left: technical
add_rect(sl, Inches(0.4), Inches(1.35), Inches(5.9), Inches(5.8),
         fill=RGBColor(0xF0, 0xF4, 0xFF), line=BLUE)
add_textbox(sl, Inches(0.55), Inches(1.45), Inches(5.6), Inches(0.38),
            "Technical Achievements", font_size=13, bold=True, color=BLUE)
bullet_box(sl, tech, Inches(0.55), Inches(1.85), Inches(5.6), Inches(5.1),
           icon="✅", font_size=11)

# Right: business
add_rect(sl, Inches(6.9), Inches(1.35), Inches(5.95), Inches(5.8),
         fill=RGBColor(0xD4, 0xED, 0xDA), line=GREEN)
add_textbox(sl, Inches(7.05), Inches(1.45), Inches(5.65), Inches(0.38),
            "Business Logic Achievements", font_size=13, bold=True, color=GREEN)
bullet_box(sl, biz, Inches(7.05), Inches(1.85), Inches(5.65), Inches(5.1),
           icon="✅", font_size=11)

# Middle separator
add_rect(sl, Inches(6.35), Inches(1.35), Inches(0.5), Inches(5.8), fill=WHITE)
add_textbox(sl, Inches(6.35), Inches(3.8), Inches(0.5), Inches(1.0),
            "vs", font_size=18, bold=True, color=GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — V1 vs V2 COMPARISON
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "V1 vs V2 Pipeline Comparison",
             "Same tools and data layer — different reasoning architectures")

table_slide(sl,
    headers=["Dimension", "V1 — Agentic Loop", "V2 — LangGraph"],
    rows=[
        ["Architecture", "Single agent, all 17 tools", "12 specialist nodes"],
        ["Reasoning", "Autonomous tool selection", "Structured routing by intent"],
        ["Auditability", "Tool call list", "Node-by-node execution trace"],
        ["Latency", "Faster for simple queries", "Better for complex chains"],
        ["Configurability", "Max iterations slider (3–25)", "Automatic per intent"],
        ["Best for", "Exploratory / conversational", "Repeatable / structured"],
        ["Context management", "History summarization at 12 msgs", "State graph accumulation"],
        ["Transparency", "Tool names + input preview", "Full node output map"],
    ],
    top=Inches(1.35),
    row_height=Inches(0.69),
)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — IMPROVEMENTS
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "10 Improvements: Decision Support → Decision Execution")

table_slide(sl,
    headers=["#", "Capability", "Business Impact"],
    rows=[
        ["1", "Decision Memory & Learning", "Recalibrate elasticity from actual outcomes"],
        ["2", "Alert Triage Queue", "Proactively surface top 5 revenue-at-risk decisions daily"],
        ["3", "Cross-Category Elasticity", "Huggies price up → store-brand diapers up 9%"],
        ["4", "Probabilistic Outcome Trees", "40% revenue up / 35% flat / 25% down"],
        ["5", "NL ERP Execution", "One-click: AI recommendation → SAP PO creation"],
        ["6", "Multi-Horizon Planning", "2W operational / 8W tactical / 26W strategic simultaneous"],
        ["7", "Stakeholder-Specific Framing", "Same analysis reframed for Buyer vs CFO vs Category Mgr"],
        ["8", "What-If Branching", "Named scenario versions with rollback capability"],
        ["9", "Competitor Move Simulation", "Costco drops 15% → AI recommends counter-strategy"],
        ["10", "Executive Summary Auto-Gen", "1-page brief auto-generated for steering committee"],
    ],
    top=Inches(1.35),
    row_height=Inches(0.575),
)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — MACRO → MICRO MENTAL MODEL
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "Macro → Micro: How Every Decision Ripples",
             "The AI reasons through all six layers simultaneously")

chain = [
    ("TRIGGER",             "Price change / Supply event / Demand shift / Tariff / Seasonal / Competitive", BLUE),
    ("DEMAND IMPACT",       "Elasticity-adjusted volume · Promotional uplift · Forecast accuracy applied", RGBColor(0xF3, 0x9C, 0x12)),
    ("INVENTORY IMPACT",    "Reorder points · Safety stock · Perishable cap · Planogram capacity check", GREEN),
    ("SUPPLY CHAIN IMPACT", "Carrier capacity · Regional coverage gaps · Lead time risk (+30% late)", RED),
    ("FINANCIAL IMPACT",    "Revenue × Margin × Vendor trade netting × VMI exclusion × Carrying cost", RGBColor(0x8E, 0x44, 0xAD)),
    ("RANKED RECOMMENDATION", "3 options: Proceed / Modify / Defer — with quantified trade-offs", BLUE),
]

box_w = Inches(12.5)
box_h = Inches(0.75)
for i, (label, detail, color) in enumerate(chain):
    top = Inches(1.38) + box_h * i
    add_rect(sl, Inches(0.4), top, Inches(2.5), box_h - Inches(0.05),
             fill=color)
    add_textbox(sl, Inches(0.5), top + Inches(0.15), Inches(2.3), box_h - Inches(0.1),
                label, font_size=11, bold=True, color=WHITE)
    add_rect(sl, Inches(2.95), top, Inches(10.0), box_h - Inches(0.05),
             fill=LIGHTGREY, line=RGBColor(0xDE, 0xE2, 0xE6))
    add_textbox(sl, Inches(3.1), top + Inches(0.15), Inches(9.7), box_h - Inches(0.1),
                detail, font_size=12, color=DARK)
    if i < len(chain) - 1:
        add_textbox(sl, Inches(1.2), top + box_h - Inches(0.1), Inches(0.5), Inches(0.25),
                    "↓", font_size=14, bold=True, color=GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — KEY NUMBERS
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)
slide_header(sl, "Key Numbers Every Specialist Should Know",
             "Hard constraints and parameters baked into every tool call")

table_slide(sl,
    headers=["Parameter", "Value", "Why It Matters"],
    rows=[
        ["Diaper price elasticity", "-1.4", "10% price hike → 14% volume drop"],
        ["Diaper recovery factor", "0.70", "Price cuts only recover 70% of lost demand"],
        ["Tobacco recovery factor", "0.40", "Very sticky — promotions don't move volume"],
        ["Dairy max days-of-supply", "3 days", "Hard cap — no exceptions regardless of signal"],
        ["VMI share (Milk)", "30%", "30% of milk inventory owned by vendor — excluded from cost"],
        ["Replenishment lag", "3–4 days", "Order today = earliest shelves in 4 days"],
        ["Late delivery probability", "30%", "30% chance of +3 extra days on top of base lag"],
        ["Forecast accuracy gate", "60%", "Below this: forecast blocked from entering PO system"],
        ["TruckCo D premium (SE diapers)", "+45%", "Only alternate carrier in SE region for diapers"],
        ["Annual carrying rate", "25%", "Cost of holding inventory for one year"],
    ],
    top=Inches(1.35),
    row_height=Inches(0.575),
)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CLOSING / DEMO
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(blank_layout)

add_rect(sl, 0, 0, SLIDE_W, SLIDE_H, fill=BLUE)
add_rect(sl, 0, Inches(5.5), SLIDE_W, Inches(2.0), fill=DARKBLUE)

add_textbox(sl, Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.0),
            "Live Demo", font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.0), Inches(2.7), Inches(11.3), Inches(0.6),
            "🌐  retail-supply-chain-ai-hp2hz8kf9cjqfkr82wogkt.streamlit.app",
            font_size=16, color=RGBColor(0xAD, 0xD8, 0xE6), align=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.0), Inches(3.5), Inches(11.3), Inches(0.5),
            "GitHub:  github.com/Krishhs89/retail-supply-chain-ai",
            font_size=14, color=RGBColor(0xAD, 0xD8, 0xE6), align=PP_ALIGN.CENTER, italic=True)

add_textbox(sl, Inches(1.0), Inches(5.65), Inches(11.3), Inches(0.5),
            "Built with  Claude claude-sonnet-4-6  ·  LangGraph  ·  Streamlit  ·  Python 3.13",
            font_size=13, color=RGBColor(0xCC, 0xE5, 0xFF), align=PP_ALIGN.CENTER)

add_textbox(sl, Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.5),
            "Interview Kickstart — Agentic AI Capstone Project  ·  2026",
            font_size=12, color=RGBColor(0x99, 0xC0, 0xD8), align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════
out = "presentation.pptx"
prs.save(out)
print(f"Saved: {out}  ({prs.slides.__len__()} slides)")
